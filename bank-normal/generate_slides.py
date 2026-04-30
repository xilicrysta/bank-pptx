import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# --- Constants & Design System (Normal / Simplified) ---
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1C, 0x1B, 0x19)
SUB_GREY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
ACCENT_LINE = RGBColor(0xD0, 0xD0, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_SANS = "Arial"

class PresentationGenerator:
    def __init__(self, title, subtitle):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle

    def _add_base_design(self, slide):
        # A simple top decorative line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_LINE
        line.line.visible = False


    def add_cover_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        if os.path.exists("assets/cover_bg.png"):
            slide.shapes.add_picture("assets/cover_bg.png", 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        else:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG_WHITE
            bg.line.visible = False
        
        # Center Box
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2), Inches(10.333), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = ACCENT_LINE
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.name = FONT_SANS
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = self.subtitle
        p2.font.name = FONT_SANS
        p2.font.size = Pt(24)
        p2.font.color.rgb = SUB_GREY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(30)

    def add_transition_slide(self, section_num, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        NAVY_BLUE = RGBColor(0x1B, 0x2B, 0x4A)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
        SUB_GREY = RGBColor(0x8A, 0x93, 0xA8)
        ACTIVE_BG = RGBColor(0xF0, 0xF4, 0xF8)

        # Background Gradient
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.line.visible = False
        
        gradient_xml = """
        <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
          <a:gsLst>
            <a:gs pos="0">
              <a:srgbClr val="1B2B4A"/>
            </a:gs>
            <a:gs pos="37500">
              <a:srgbClr val="1B2B4A"/>
            </a:gs>
            <a:gs pos="100000">
              <a:srgbClr val="465A82"/>
            </a:gs>
          </a:gsLst>
          <a:lin ang="0" scaled="1"/>
        </a:gradFill>
        """
        
        spPr = bg._element.spPr
        solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            spPr.remove(solidFill)
        spPr.append(parse_xml(gradient_xml))

        # Left bar accent
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(0.15), Inches(5.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.visible = False
        
        num_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1))
        p = num_box.text_frame.paragraphs[0]
        if section_num:
            p.text = f"CHAPTER {section_num}"
        else:
            p.text = ""
        p.font.name = FONT_SANS
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GREY
        
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(8), Inches(2))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT_SANS
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        title_box.text_frame.word_wrap = True

        chapters = [
            "導入 (振り返り)",
            "銀行の基本的な仕組み",
            "銀行の本質",
            "テクノロジーと異業種",
            "地方銀行と投資ファンド",
            "銀行のこれからの姿",
            "まとめ"
        ]
        
        if section_num:
            try:
                curr_idx = int(section_num)
            except ValueError:
                curr_idx = -1
        else:
            if "そもそも何" in title:
                curr_idx = 0
            elif "まとめ" in title:
                curr_idx = 6
            else:
                curr_idx = -1
            
        start_y = Inches(1.0)
        box_h = Inches(0.55)
        spacing = Inches(0.2)
        
        # Draw the skewer progress line
        line_x = Inches(11.15)
        line_w = Inches(0.08)
        skewer_start_y = start_y + box_h / 2
        skewer_h = 6 * (box_h + spacing)
        
        # Background track
        skewer_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x - line_w/2, skewer_start_y, line_w, skewer_h)
        skewer_bg.fill.solid()
        skewer_bg.fill.fore_color.rgb = RGBColor(0x2B, 0x3B, 0x5A)
        skewer_bg.line.visible = False
        
        # Progress track
        if curr_idx > 0:
            prog_h = curr_idx * (box_h + spacing)
            prog_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x - line_w/2, skewer_start_y, line_w, prog_h)
            prog_line.fill.solid()
            prog_line.fill.fore_color.rgb = WHITE
            prog_line.line.visible = False
        
        for i, chap in enumerate(chapters):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(9.5), 
                start_y + i * (box_h + spacing), 
                Inches(3.3), 
                box_h
            )
            box.fill.solid()
            if i == curr_idx:
                # 現在アクティブ
                box.fill.fore_color.rgb = ACTIVE_BG
                box.line.visible = False
                font_color = NAVY_BLUE
            elif i < curr_idx:
                # 完了済み
                box.fill.fore_color.rgb = NAVY_BLUE
                box.line.color.rgb = WHITE
                box.line.width = Pt(0.75)
                font_color = SUB_GREY
            else:
                # 未完了
                box.fill.fore_color.rgb = RGBColor(0x2B, 0x3B, 0x5A)
                box.line.visible = False
                font_color = WHITE
            
            tf = box.text_frame
            tf.word_wrap = True
            p_c = tf.paragraphs[0]
            p_c.text = chap
            p_c.font.name = FONT_SANS
            p_c.font.size = Pt(14)
            p_c.font.bold = True
            p_c.font.color.rgb = font_color
            
            if i < len(chapters) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    Inches(11.05),
                    start_y + i * (box_h + spacing) + box_h + Inches(0.05),
                    Inches(0.2),
                    Inches(0.15)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = SUB_GREY
                arrow.line.visible = False

    def add_message_slide(self, title, message, body_items=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # Force white bg
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.visible = False

        self._add_base_design(slide)
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(1))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT_SANS
        p.font.size = Pt(24)
        p.font.color.rgb = SUB_GREY
        
        # Message
        m_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2))
        tf = m_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = message
        p.font.name = FONT_SANS
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        if body_items:
            b_box = slide.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(12.5), Inches(3.8))
            tf = b_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            for item in body_items:
                p = tf.add_paragraph()
                p.text = f"● {item}"
                p.font.name = FONT_SANS
                p.font.size = Pt(26)
                p.space_after = Pt(24)
                p.font.color.rgb = TEXT_DARK

    def add_diagram_slide(self, title, message, diagram_func):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.visible = False

        self._add_base_design(slide)
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(1))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT_SANS
        p.font.size = Pt(24)
        p.font.color.rgb = SUB_GREY
        
        # Message
        m_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2))
        tf = m_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = message
        p.font.name = FONT_SANS
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        diagram_func(slide)

    # Simplified diagrams
    def draw_megabank_flow(self, slide):
        # Lots of banks -> Arrow -> 3 Megabanks
        y = Inches(3.5)
        
        # Left box
        b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(3.5), Inches(1.5))
        b1.fill.solid()
        b1.fill.fore_color.rgb = LIGHT_GREY
        b1.line.color.rgb = ACCENT_LINE
        p = b1.text_frame.paragraphs[0]
        p.text = "昔の日本\n(たくさんの銀行があった)"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        b1.text_frame.word_wrap = True
        
        # Arrow
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.5), y+Inches(0.2), Inches(2.3), Inches(1.1))
        arr.fill.solid()
        arr.fill.fore_color.rgb = TEXT_DARK
        arr.line.visible = False
        arr.text_frame.paragraphs[0].text = "お互いに合体！\n(合併)"
        arr.text_frame.paragraphs[0].font.size = Pt(16)
        arr.text_frame.paragraphs[0].font.color.rgb = WHITE
        
        # Right box
        b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), y-Inches(0.5), Inches(3.5), Inches(2.5))
        b2.fill.solid()
        b2.fill.fore_color.rgb = TEXT_DARK
        b2.line.visible = False
        p = b2.text_frame.paragraphs[0]
        p.text = "今の3大メガバンク\n\n・三菱UFJ銀行\n・三井住友銀行\n・みずほ銀行"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        b2.text_frame.word_wrap = True

    def draw_credit_creation(self, slide):
        y = Inches(4.0)
        
        shape1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(2), Inches(2))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = LIGHT_GREY
        shape1.line.color.rgb = ACCENT_LINE
        shape1.text_frame.paragraphs[0].text = "Aさんの預金\n(100円)"
        shape1.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        shape1.text_frame.paragraphs[0].font.bold = True
        
        slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), y+Inches(0.8), Inches(0.8), Inches(0.4)).fill.solid()
        
        shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), y, Inches(2), Inches(2))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = LIGHT_GREY
        shape2.line.color.rgb = ACCENT_LINE
        shape2.text_frame.paragraphs[0].text = "Bさんに貸出\n(90円)"
        shape2.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        shape2.text_frame.paragraphs[0].font.bold = True

        slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.4), y+Inches(0.8), Inches(0.8), Inches(0.4)).fill.solid()

        shape3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), y, Inches(2), Inches(2))
        shape3.fill.solid()
        shape3.fill.fore_color.rgb = LIGHT_GREY
        shape3.line.color.rgb = ACCENT_LINE
        shape3.text_frame.paragraphs[0].text = "別の預金へ\n(90円)"
        shape3.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        shape3.text_frame.paragraphs[0].font.bold = True
        
        slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.6), y+Inches(0.8), Inches(0.8), Inches(0.4)).fill.solid()
        
        shape4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), y-Inches(0.25), Inches(2.5), Inches(2.5))
        shape4.fill.solid()
        shape4.fill.fore_color.rgb = TEXT_DARK
        shape4.line.visible = False
        shape4.text_frame.paragraphs[0].text = "世の中のお金\n= 190円に！"
        shape4.text_frame.paragraphs[0].font.color.rgb = WHITE
        shape4.text_frame.paragraphs[0].font.size = Pt(24)
        shape4.text_frame.paragraphs[0].font.bold = True
        
        lbl = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11), Inches(1))
        lbl.text_frame.paragraphs[0].text = "銀行を挟んで貸し借りを繰り返すだけで、世の中のお金が増える不思議な仕組み。"
        lbl.text_frame.paragraphs[0].font.size = Pt(22)

    def draw_interest_margin(self, slide):
        base_x = Inches(1.0)
        y = Inches(3.5)
        
        # Deposit
        b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base_x, y, Inches(3.3), Inches(2.3))
        b1.fill.solid()
        b1.fill.fore_color.rgb = LIGHT_GREY
        b1.line.color.rgb = ACCENT_LINE
        p1 = b1.text_frame.paragraphs[0]
        p1.text = "① みんなから預かる\n(利息 1円を払う)"
        p1.font.color.rgb = TEXT_DARK
        p1.font.size = Pt(22)
        p1.font.bold = True
        
        # Loan
        b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base_x + Inches(3.8), y, Inches(3.3), Inches(2.3))
        b2.fill.solid()
        b2.fill.fore_color.rgb = SUB_GREY
        b2.line.visible = False
        p2 = b2.text_frame.paragraphs[0]
        p2.text = "② 企業に貸す\n(利息 10円をもらう)"
        p2.font.color.rgb = WHITE
        p2.font.size = Pt(22)
        p2.font.bold = True
        
        # Margin
        b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base_x + Inches(7.6), y, Inches(3.5), Inches(2.3))
        b3.fill.solid()
        b3.fill.fore_color.rgb = TEXT_DARK
        b3.line.visible = False
        p3 = b3.text_frame.paragraphs[0]
        p3.text = "③ 銀行の利益\n(差額の 9円)"
        p3.font.color.rgb = WHITE
        p3.font.size = Pt(26)
        p3.font.bold = True
        
        lbl = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(1))
        lbl.text_frame.paragraphs[0].text = "スーパーの「安く仕入れて、高く売る」と全く同じ『利ざや』の考え方である。"
        lbl.text_frame.paragraphs[0].font.size = Pt(24)


    def draw_balance_sheet(self, slide):
        y = Inches(3.5)
        w = Inches(4.5)
        h = Inches(3.0)
        center = self.prs.slide_width / 2
        
        left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center - w - Inches(0.2), y, w, h)
        left.fill.solid()
        left.fill.fore_color.rgb = LIGHT_GREY
        left.line.color.rgb = ACCENT_LINE
        left.text_frame.paragraphs[0].text = "■ お金の使い道 (資産)\n\n企業に貸しているお金など"
        left.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        left.text_frame.paragraphs[0].font.size = Pt(22)
        left.text_frame.paragraphs[0].font.bold = True
        left.text_frame.word_wrap = True
        
        right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center + Inches(0.2), y, w, h)
        right.fill.solid()
        right.fill.fore_color.rgb = TEXT_DARK
        right.line.visible = False
        right.text_frame.paragraphs[0].text = "■ お金の集め方 (負債)\n\nみんなから預かっているお金"
        right.text_frame.paragraphs[0].font.color.rgb = WHITE
        right.text_frame.paragraphs[0].font.size = Pt(22)
        right.text_frame.paragraphs[0].font.bold = True
        right.text_frame.word_wrap = True
        
        lbl = slide.shapes.add_textbox(center - Inches(3), y - Inches(0.8), Inches(6), Inches(0.5))
        lbl.text_frame.paragraphs[0].text = "私たちが見ると逆！みんなの「預金」は銀行にとって「借金」"
        lbl.text_frame.paragraphs[0].font.size = Pt(20)
        lbl.text_frame.paragraphs[0].font.bold = True
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def draw_industry_map(self, slide):
        groups = [
            ("日本銀行", "全てのお金の元締め (紙幣を刷る)", Inches(0.6), Inches(3.0)),
            ("メガバンク", "全国展開・世界相手 (三菱UFJなど)", Inches(3.7), Inches(3.0)),
            ("地方銀行", "あなたの地元の企業を応援している", Inches(6.8), Inches(3.0)),
            ("ネット銀行", "スマホで完結・店舗がない (楽天など)", Inches(9.9), Inches(3.0))
        ]
        for name, desc, x, y in groups:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(2.5))
            box.fill.solid()
            if name == "日本銀行":
                box.fill.fore_color.rgb = TEXT_DARK
                box.text_frame.paragraphs[0].font.color.rgb = WHITE
            else:
                box.fill.fore_color.rgb = LIGHT_GREY
                box.line.color.rgb = ACCENT_LINE
                box.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
                
            p = box.text_frame.paragraphs[0]
            p.text = f"{name}\n\n{desc}"
            p.font.size = Pt(18)
            p.font.bold = True
            box.text_frame.word_wrap = True

    def add_summary_slide(self, items):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.visible = False

        self._add_base_design(slide)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = "今日のまとめ"
        p.font.name = FONT_SANS
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        for item in items:
            p = tf.add_paragraph()
            p.text = f"✔ {item}"
            p.font.name = FONT_SANS
            p.font.size = Pt(28)
            p.space_after = Pt(24)
            p.font.color.rgb = TEXT_DARK

    def save(self, filename):
        # 確実に日本語フォントが適用されるよう、a:ea (East Asian) を全テキストに設定
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Arial"
                    try:
                        defRPr = paragraph.font._element
                        defRPr.set(qn('a:ea'), "Meiryo")
                    except Exception:
                        pass
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        try:
                            rPr = run.font._element
                            rPr.set(qn('a:ea'), "Meiryo")
                        except Exception:
                            pass
                            
        self.prs.save(filename)


    def draw_stablecoin(self, slide):
        y = Inches(3.5)
        # Digital Coin
        b1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), y, Inches(2.0), Inches(2.0))
        b1.fill.solid()
        b1.fill.fore_color.rgb = LIGHT_GREY
        b1.line.color.rgb = ACCENT_LINE
        p1 = b1.text_frame.paragraphs[0]
        p1.text = "デジタル通貨\n(仮想通貨など)"
        p1.font.color.rgb = TEXT_DARK
        p1.font.size = Pt(20)
        p1.font.bold = True
        
        # Link
        slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(4.5), y+Inches(0.8), Inches(1.5), Inches(0.4)).fill.solid()
        
        # Fiat
        b2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), y, Inches(2.0), Inches(2.0))
        b2.fill.solid()
        b2.fill.fore_color.rgb = LIGHT_GREY
        b2.line.color.rgb = ACCENT_LINE
        p2 = b2.text_frame.paragraphs[0]
        p2.text = "法定通貨\n(円やドル)"
        p2.font.color.rgb = TEXT_DARK
        p2.font.size = Pt(20)
        p2.font.bold = True
        
        lbl = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11), Inches(1))
        lbl.text_frame.paragraphs[0].text = "価格が「円」や「ドル」と連動（固定）しているため、支払いや送金に使いやすい。"
        lbl.text_frame.paragraphs[0].font.size = Pt(24)

    def draw_smartphone_payment(self, slide):
        y = Inches(4.5)
        # User
        b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(2.5), Inches(1.5))
        b1.fill.solid()
        b1.fill.fore_color.rgb = LIGHT_GREY
        b1.line.color.rgb = ACCENT_LINE
        p1 = b1.text_frame.paragraphs[0]
        p1.text = "私たち\n(スマホ決済)"
        p1.font.color.rgb = TEXT_DARK
        p1.font.size = Pt(22)
        p1.font.bold = True
        
        slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.0), y+Inches(0.5), Inches(1.0), Inches(0.4)).fill.solid()
        
        # Store
        b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), y, Inches(2.5), Inches(1.5))
        b2.fill.solid()
        b2.fill.fore_color.rgb = LIGHT_GREY
        b2.line.color.rgb = ACCENT_LINE
        p2 = b2.text_frame.paragraphs[0]
        p2.text = "お店"
        p2.font.color.rgb = TEXT_DARK
        p2.font.size = Pt(22)
        p2.font.bold = True

        slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(2.0), y-Inches(1.0), Inches(0.4), Inches(0.8)).fill.solid()
        
        # Bank
        b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y-Inches(1.5), Inches(7.0), Inches(1.0))
        b3.fill.solid()
        b3.fill.fore_color.rgb = TEXT_DARK
        b3.line.visible = False
        p3 = b3.text_frame.paragraphs[0]
        p3.text = "銀行口座 (裏側でチャージや精算を行っている)"
        p3.font.color.rgb = WHITE
        p3.font.size = Pt(22)
        p3.font.bold = True

    def draw_pe_vc_flow(self, slide):
        y = Inches(3.0)
        # Bank
        b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(2.5), Inches(2.0))
        b1.fill.solid()
        b1.fill.fore_color.rgb = LIGHT_GREY
        b1.line.color.rgb = ACCENT_LINE
        p1 = b1.text_frame.paragraphs[0]
        p1.text = "銀行\n(融資)"
        p1.font.color.rgb = TEXT_DARK
        p1.font.size = Pt(22)
        p1.font.bold = True
        
        # VC/PE
        b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), y, Inches(3.0), Inches(2.0))
        b2.fill.solid()
        b2.fill.fore_color.rgb = TEXT_DARK
        b2.line.visible = False
        p2 = b2.text_frame.paragraphs[0]
        p2.text = "投資ファンド\n(PE / VC)\n(株式出資・経営支援)"
        p2.font.color.rgb = WHITE
        p2.font.size = Pt(22)
        p2.font.bold = True

        slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.0), y+Inches(0.8), Inches(1.0), Inches(0.4)).fill.solid()

        # Company
        b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), y, Inches(3.0), Inches(2.0))
        b3.fill.solid()
        b3.fill.fore_color.rgb = LIGHT_GREY
        b3.line.color.rgb = ACCENT_LINE
        p3 = b3.text_frame.paragraphs[0]
        p3.text = "成長企業・\n後継者不足の企業"
        p3.font.color.rgb = TEXT_DARK
        p3.font.size = Pt(22)
        p3.font.bold = True
        
        lbl = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.333), Inches(1))
        lbl.text_frame.paragraphs[0].text = "銀行がお金を貸すだけでなく、ファンドと協力して企業の成長や再生を直接支援するケースが増えている。"
        lbl.text_frame.paragraphs[0].font.size = Pt(22)

# --- Generator Script ---
def generate_normal_presentation():
    gen = PresentationGenerator("みんなの銀行基礎知識", "〜 難しくない！お金と銀行のリアルな話 〜")
    
    # 1. 導入 & 基本的な仕組み
    gen.add_cover_slide()
    gen.add_message_slide("導入", "銀行とは、何をするところか", [
        "お金を安全に預けておく「巨大な金庫」",
        "毎月の給料が振り込まれる「自分の口座」",
        "実は、ただの「お金の置き場所」ではない",
        "社会中にお金をくまなく巡らせる『心臓』としての役割がある"
    ])
    gen.add_transition_slide("", "銀行ってそもそも何なのか")
    gen.add_message_slide("本日のアジェンダ", "今日話すこと", [
        "1. 銀行の基本的な仕組み",
        "2. 銀行の本質（なぜ存在するのか）",
        "3. 新しいテクノロジーと異業種の参入",
        "4. 地方銀行の現状と投資ファンドの役割",
        "5. 銀行のこれからの姿"
    ])
    
    gen.add_transition_slide("1", "銀行の基本的な仕組み")
    gen.add_message_slide("銀行の三大業務（大切な仕事）", "銀行が絶対にやっている3つのこと", [
        "預金 (よきん) ： みんなのお金を安全に預かる",
        "融資 (ゆうし) ： お金に困っている人や頑張る企業に貸す",
        "為替 (かわせ) ： 遠くの人へ、現金を直接運ばずにお金を送金する"
    ])
    gen.add_message_slide("お金のめぐり方", "銀行は「お金の橋渡し」", [
        "お金が余っている人（私たち）から預かる。",
        "お金が足りない人や企業に貸す。",
        "銀行という橋があるから、社会全体でお金がスムーズに回る。"
    ])
    gen.add_diagram_slide("魔法の仕組み：信用創造", "銀行を挟むと、なぜかお金の総額が増える仕組み", gen.draw_credit_creation)
    gen.add_message_slide("信用創造のタネあかし", "「データ上のお金」が増えているだけ", [
        "実際に現金の札束や硬貨が増えているわけではない。",
        "通帳の残高に「○○円ある」と記録されることで、お金として使える。",
        "この仕組みのおかげで、企業は事業を大きくすることができる。"
    ])
    gen.add_diagram_slide("どうやって儲けるか (利ざや)", "普通の商売と同じ「安く仕入れて、高く売る」", gen.draw_interest_margin)
    gen.add_message_slide("他の儲け方 (手数料ビジネス)", "利ざや以外にもこんな稼ぎ方がある", [
        "振込手数料：お金を送る時にかかる「送料」のようなもの",
        "ATM手数料：時間外や他行の引き出しなどで発生するお金",
        "投資信託・保険：他の金融商品を「代わりに販売する」ことで貰える報酬"
    ])
    gen.add_diagram_slide("銀行のバランスシート (持ち物リスト)", "私たちのお金は、銀行にとって『預かりもの』", gen.draw_balance_sheet)

    # 2. 銀行の本質
    gen.add_transition_slide("2", "銀行の本質（なぜ存在するのか）")
    gen.add_message_slide("銀行が存在する最大の理由", "「信用」を仲介している", [
        "見ず知らずの人に直接お金を貸すのは、返ってこないリスクが大きすぎる。",
        "銀行が間に立ち「確実に返ってくる」という信用を担保している。",
        "預金者も「銀行なら安全だ」と信用してお金を預けることができる。"
    ])
    gen.add_diagram_slide("銀行の歴史", "昔はたくさんあった銀行が、合体して巨大化した。", gen.draw_megabank_flow)
    
    # 3. テクノロジーと異業種参入
    gen.add_transition_slide("3", "新しいテクノロジーと異業種の参入")
    gen.add_diagram_slide("スマホ決済の台頭 (PayPayやメルペイ)", "銀行に行かなくても支払いができる世界へ", gen.draw_smartphone_payment)
    gen.add_message_slide("ネット銀行の急成長", "「スマホの中にある便利な銀行」", [
        "店舗を持たないため、家賃や人件費がかからず、手数料が安い。",
        "ポイントが貯まりやすかったり、アプリが使いやすかったりする。"
    ])
    gen.add_diagram_slide("デジタル通貨の波：ステーブルコイン", "ブロックチェーン技術を使った新しいお金の形", gen.draw_stablecoin)
    gen.add_message_slide("テクノロジーがもたらす変化", "金融サービスがもっと身近に", [
        "銀行以外のIT企業が、次々と金融サービスを提供し始めている。",
        "送金コストが限りなくゼロに近づく未来が予想される。"
    ])

    # 4. 地銀合併とPE/VC
    gen.add_transition_slide("4", "地方銀行の現状と投資ファンドの役割")
    gen.add_message_slide("今の銀行の悩み：金利がほぼ0円", "「利ざや」が稼げなくて大ピンチ", [
        "今は金利が低いため、企業に貸しても「利息」がほとんど貰えない。",
        "100万円を1年預けても、チロルチョコ1個買えない時代である。"
    ])
    gen.add_message_slide("地方銀行の「合体」と「統廃合」", "生き残りをかけて規模を大きくする", [
        "稼ぎにくくなったため、 　　地方の銀行同士が協力・合体するニュースが増えている。",
        "また、人が来ない店舗（支店）を減らして、 　　コストを下げる努力をしている。"
    ])
    gen.add_message_slide("信用金庫の役割", "利益第一ではない、地域密着のサポーター", [
        "地方銀行と似ているが、株式会社ではないため「株主の利益」を優先しない。",
        "地域の人々や中小企業が互いに助け合う「相互扶助」を目的に作られている。",
        "だからこそ、地元の小さな企業にも親身になって寄り添える。"
    ])
    gen.add_diagram_slide("投資ファンド（PE・VC）との連携", "貸すだけではなく、直接経営を支援する時代へ", gen.draw_pe_vc_flow)

    # 5. これからの姿
    gen.add_transition_slide("5", "銀行のこれからの姿")
    gen.add_message_slide("DX (デジタル化) の推進", "店舗のデジタル完結へ", [
        "「店舗に行って、書類を書いて待つ」という昔のやり方を変える。",
        "スマホアプリひとつで全部完結する「超ベンリな銀行」へ進化中である。"
    ])
    gen.add_message_slide("これからのビジネスモデル", "「お金を貸すだけ」からの卒業", [
        "お金の貸し借りだけでなく、みんなの「生活の悩み」を解決するアドバイザーへ。",
        "他の便利なアプリの「裏側」に、銀行の仕組み・機能だけをこっそり提供する。"
    ])

    # まとめ
    gen.add_transition_slide("", "まとめ")
    gen.add_summary_slide([
        "銀行は「信用」を仲介し、社会にお金を巡らせる「心臓」の役割を果たす。",
        "今は金利が低く、地方銀行は合併やファンドとの協力で生き残りを図っている。",
        "テクノロジーの進化により、IT企業など異業種との境界線がなくなりつつある。",
        "これからはスマホの奥に溶け込む「見えない銀行」へと進化していく。"
    ])

    output_path = "bank_normal_presentation.pptx"
    gen.save(output_path)
    print(f"Normal Presentation saved to {output_path}")

if __name__ == "__main__":
    generate_normal_presentation()
