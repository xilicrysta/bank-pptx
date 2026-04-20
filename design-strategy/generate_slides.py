import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Constants & Design System (Ver 2. Strategy) ---
BG_WHITE = RGBColor(0xFE, 0xFE, 0xFE)
TEXT_DARK = RGBColor(0x1C, 0x1B, 0x19)
ACCENT_YELLOW = RGBColor(0xF9, 0xD1, 0x3E)
SUB_GREY = RGBColor(0x82, 0x82, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_SANS = "Hiragino Sans"
FONT_SANS_BOLD = "Hiragino Sans"

class PresentationGenerator:
    def __init__(self, title, subtitle):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle

    def _add_frame(self, slide, color=TEXT_DARK):
        """Adds a thick rounded rectangle frame characteristic of the design."""
        frame_margin = Inches(0.2)
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            frame_margin, frame_margin, 
            self.prs.slide_width - (frame_margin * 2), 
            self.prs.slide_height - (frame_margin * 2)
        )
        frame.fill.background()
        line = frame.line
        line.color.rgb = color
        line.width = Pt(4)
        frame.adjustments[0] = 0.02

    def _add_pill_label(self, slide, text, x, y, bg_color=ACCENT_YELLOW, text_color=TEXT_DARK):
        """Adds a rounded 'pill' style label."""
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.5), Inches(0.4))
        pill.fill.solid()
        pill.fill.fore_color.rgb = bg_color
        pill.line.visible = False
        pill.adjustments[0] = 0.5 
        
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_SANS
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    def add_cover_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_frame(slide)
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.name = FONT_SANS
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = self.subtitle
        p2.font.name = FONT_SANS
        p2.font.size = Pt(24)
        p2.font.color.rgb = SUB_GREY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(30)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.5), Inches(2.3), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_YELLOW
        bar.line.visible = False

    def add_transition_slide(self, section_num, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_frame(slide, color=ACCENT_YELLOW)
        
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
        tf_n = num_box.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = f"Section {section_num:02d}"
        p_n.font.name = FONT_SANS
        p_n.font.size = Pt(60)
        p_n.font.bold = True
        p_n.font.color.rgb = ACCENT_YELLOW
        p_n.alignment = PP_ALIGN.CENTER
        
        display_title = title
        if ". " in title:
            parts = title.split(". ", 1)
            if parts[0].isdigit():
                display_title = parts[1]

        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1.5))
        tf_t = title_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = display_title
        p_t.font.name = FONT_SANS
        p_t.font.size = Pt(40)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK
        p_t.alignment = PP_ALIGN.CENTER

    def add_message_slide(self, title, message, body_items=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_frame(slide)
        self._add_pill_label(slide, title, Inches(0.5), Inches(0.5))
        
        msg_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(1.2))
        tf_m = msg_box.text_frame
        p_m = tf_m.paragraphs[0]
        p_m.text = message
        p_m.font.name = FONT_SANS
        p_m.font.size = Pt(36)
        p_m.font.bold = True
        p_m.font.color.rgb = TEXT_DARK
        p_m.alignment = PP_ALIGN.LEFT
        
        if body_items:
            # Shifted as far left as possible and widened to the limit to prevent any wrapping issues
            body_box = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(12.5), Inches(3.5))
            tf_b = body_box.text_frame
            tf_b.word_wrap = True
            tf_b.vertical_anchor = MSO_ANCHOR.TOP
            for item in body_items:
                p = tf_b.add_paragraph()
                p.text = f"• {item}"
                p.font.name = FONT_SANS
                p.font.size = Pt(23) # Increased from 21 to improve readability
                p.space_after = Pt(16)
                p.font.color.rgb = TEXT_DARK

    def add_diagram_slide(self, title, message, diagram_func):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_frame(slide)
        self._add_pill_label(slide, title, Inches(0.5), Inches(0.5))
        
        msg_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.8))
        p_m = msg_box.text_frame.paragraphs[0]
        p_m.text = message
        p_m.font.name = FONT_SANS
        p_m.font.size = Pt(28)
        p_m.font.bold = True
        p_m.font.color.rgb = TEXT_DARK
        
        diagram_func(slide)

    def draw_megabank_flow(self, slide):
        y_start = Inches(3.0)
        box_w = Inches(3.5)
        box_h = Inches(0.8)
        
        headers = ["旧財閥・大型店", "合併・統合の流れ", "現3大メガバンク"]
        for i, h in enumerate(headers):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5 + i*4.0), y_start - Inches(1), Inches(3.0), Inches(0.5))
            box.fill.solid()
            box.fill.fore_color.rgb = SUB_GREY
            box.line.visible = False
            box.text_frame.paragraphs[0].text = h
            box.text_frame.paragraphs[0].font.size = Pt(14)
            box.text_frame.paragraphs[0].font.color.rgb = WHITE

        groups = [
            ("第一勧銀・富士・日本興業", "みずほFG", RGBColor(0x30, 0x3F, 0x9F)),
            ("三菱・東京・UFJ", "三菱UFJ FG", RGBColor(0xD3, 0x2F, 0x2F)),
            ("住友・さくら", "三井住友 FG", RGBColor(0x38, 0x8E, 0x3C))
        ]
        for i, (old, new, color) in enumerate(groups):
            curr_y = y_start + i * Inches(1.2)
            # Old
            b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), curr_y, box_w, box_h)
            b1.fill.solid()
            b1.fill.fore_color.rgb = SUB_GREY
            b1.line.visible = False
            b1.text_frame.paragraphs[0].text = old
            b1.text_frame.paragraphs[0].font.size = Pt(16)
            # Arrow
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), curr_y + Inches(0.2), Inches(1.0), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_YELLOW
            arrow.line.visible = False
            # New
            b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), curr_y, box_w, box_h)
            b2.fill.solid()
            b2.fill.fore_color.rgb = color
            b2.line.visible = False
            b2.text_frame.paragraphs[0].text = new
            b2.text_frame.paragraphs[0].font.size = Pt(20)
            b2.text_frame.paragraphs[0].font.bold = True
            b2.text_frame.paragraphs[0].font.color.rgb = WHITE

    def draw_credit_creation_diagram(self, slide):
        y = Inches(3.0)
        size = Inches(2.2)
        gap = Inches(0.4)
        steps = ["A銀行への預金\n(100万円)", "B社への貸出\n(90万円)", "C銀行への預金\n(90万円)", "世の中のお金\n= 190万円"]
        for i, s in enumerate(steps):
            x = Inches(1.0 + i * 2.6) # 2.6 = 2.2 size + 0.4 gap
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE if i < 3 else MSO_SHAPE.OVAL, x, y, size, size)
            box.fill.solid()
            box.fill.fore_color.rgb = ACCENT_YELLOW if i == 3 else SUB_GREY
            box.line.color.rgb = TEXT_DARK
            box.line.width = Pt(2)
            box.text_frame.paragraphs[0].text = s
            box.text_frame.paragraphs[0].font.size = Pt(18)
            box.text_frame.paragraphs[0].font.bold = True
            box.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
            if i < 3:
                arr_x = x + size
                arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_x, y + (size/2) - Inches(0.2), gap, Inches(0.4))
                arr.fill.solid()
                arr.fill.fore_color.rgb = TEXT_DARK
                arr.line.visible = False

    def draw_interest_margin_diagram(self, slide):
        base_x = Inches(4.0)
        base_y = Inches(6.0)
        width = Inches(2.0)
        # Yield
        y1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, base_x, base_y - Inches(4.0), width, Inches(4.0))
        y1.fill.solid()
        y1.fill.fore_color.rgb = ACCENT_YELLOW # Changed to yellow for visibility
        y1.text_frame.paragraphs[0].text = "貸出利回り"
        y1.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        y1.text_frame.paragraphs[0].font.bold = True
        
        # Cost Box (Explicit box above Deposit Interest - Now Light Blue and Aligned)
        cost_top = base_y - Inches(4.0)
        c_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, base_x + width + Inches(0.5), cost_top, width, Inches(3.0))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = RGBColor(0xAD, 0xD8, 0xE6) # Light Blue
        c_box.line.color.rgb = TEXT_DARK
        c_box.line.width = Pt(1)
        c_box.text_frame.paragraphs[0].text = "銀行の収益 (利ざや)"
        c_box.text_frame.paragraphs[0].font.size = Pt(14)
        c_box.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        c_box.text_frame.paragraphs[0].font.bold = True

        # Cost (Deposit Interest)
        y2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, base_x + width + Inches(0.5), base_y - Inches(1.0), width, Inches(1.0))
        y2.fill.solid()
        y2.fill.fore_color.rgb = SUB_GREY
        y2.text_frame.paragraphs[0].text = "預金利息"
        # Brace
        bracket_x = base_x + (width * 2) + Inches(1.0)
        bracket = slide.shapes.add_shape(MSO_SHAPE.RIGHT_BRACE, bracket_x, base_y - Inches(4.0), Inches(0.5), Inches(3.0))
        # Label
        label = slide.shapes.add_textbox(bracket_x + Inches(0.6), base_y - Inches(2.5), Inches(3.0), Inches(1.0))
        p = label.text_frame.paragraphs[0]
        p.text = "この差が「利ざや」\n(収益の源泉)"
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    def draw_bs_comparison(self, slide):
        y = Inches(2.5)
        w = Inches(5.0)
        h = Inches(4.0)
        center_x = self.prs.slide_width / 2
        b_x = center_x - w - Inches(0.5)
        
        slide.shapes.add_textbox(b_x, y - Inches(0.5), w, Inches(0.5)).text_frame.paragraphs[0].text = "銀行のB/S"
        # Assets
        al = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, b_x, y, w/2, h)
        al.text_frame.paragraphs[0].text = "資産\n\n(貸付金など)"
        al.fill.solid()
        al.fill.fore_color.rgb = ACCENT_YELLOW
        # Liab
        rl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, b_x + w/2, y, w/2, h*0.9)
        rl.text_frame.paragraphs[0].text = "負債\n\n(預金など)"
        rl.fill.solid()
        rl.fill.fore_color.rgb = SUB_GREY
        # Equity
        el = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, b_x + w/2, y + h*0.9, w/2, h*0.1)
        el.text_frame.paragraphs[0].text = "純資産"
        el.fill.solid()
        el.fill.fore_color.rgb = TEXT_DARK
        # Red Label (Moved to bottom to avoid overlap)
        lbl = slide.shapes.add_textbox(b_x, y + h + Inches(0.2), Inches(5.0), Inches(0.5))
        p = lbl.text_frame.paragraphs[0]
        p.text = "★銀行にとって、預金は「預かりもの」であり負債となる"
        p.font.color.rgb = RGBColor(0xFF, 0, 0)
        p.font.size = Pt(16)
        p.font.bold = True

    def draw_industry_map(self, slide):
        layers = [
            ("中央銀行", "日本銀行 (唯一無二)", TEXT_DARK),
            ("都市銀行", "メガバンク (全国・海外展開)", RGBColor(0x15, 0x65, 0xC0)),
            ("地方銀行", "地域密着型 (地元経済の要)", RGBColor(0x2E, 0x7D, 0x32)),
            ("ネット銀行", "利便性・低コスト追求", RGBColor(0xEF, 0x6C, 0x00))
        ]
        for i, (title, desc, color) in enumerate(layers):
            y = Inches(2.2 + i * 1.2)
            lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(3.0), Inches(0.8))
            lbl.fill.solid()
            lbl.fill.fore_color.rgb = color
            lbl.text_frame.paragraphs[0].text = title
            lbl.text_frame.paragraphs[0].font.bold = True
            lbl.text_frame.paragraphs[0].font.color.rgb = WHITE
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), y, Inches(7.5), Inches(0.8))
            box.fill.background()
            box.line.color.rgb = SUB_GREY
            tf_desc = box.text_frame
            p_desc = tf_desc.paragraphs[0]
            p_desc.text = desc
            p_desc.font.color.rgb = TEXT_DARK # Clear visibility
            p_desc.alignment = PP_ALIGN.LEFT

    def add_summary_slide(self, items):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_frame(slide, color=ACCENT_YELLOW)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = "SUMMARY"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        # Maximized width and shifted left to ensure long sentences fit without hanging punctation
        body_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.5), Inches(12.5), Inches(4))
        tf = body_box.text_frame
        tf.word_wrap = True
        for item in items:
            p = tf.add_paragraph()
            p.text = f"✔ {item}"
            p.font.size = Pt(23) # Slightly reduced to ensure single-line fits
            p.space_after = Pt(20)
            p.font.color.rgb = TEXT_DARK

    def save(self, filename):
        self.prs.save(filename)

def generate_strategic_presentation():
    sections = [
        "銀行の歴史（合併・再編）",
        "銀行の新たなビジネスモデル",
        "銀行業界の現状",
        "競合・新たな脅威",
        "銀行の将来・今後の展望"
    ]
    gen = PresentationGenerator("銀行業界の基礎知識", "〜 戦略・ビジネスモデル・未来への展望 〜")
    gen.add_cover_slide()
    gen.add_message_slide("導入：定義", "銀行とは、一言で言えば\n「信用というインフラを管理する装置」である。", 
                        ["余剰資金（預金）を必要箇所（融資）へ分配する。", "社会の決済システム（為替）を維持する。", "経済活動の根幹を支える『信頼の仲介者』である。"])
    gen.add_message_slide("導入：アジェンダ", "本日の論点：銀行ビジネスの本質と未来", 
                        [f"Section 01: {sections[0]}", f"Section 02: {sections[1]}", f"Section 03: {sections[2]}", f"Section 04: {sections[3]}", f"Section 05: {sections[4]}"])
    gen.add_transition_slide(1, sections[0])
    gen.add_diagram_slide("Section 01", "激動の平成：13の都市銀行が『3大メガバンク』へ集約", gen.draw_megabank_flow)
    gen.add_transition_slide(2, sections[1])
    gen.add_message_slide("Section 02", "三大業務：預金・融資・為替の相乗効果", 
                         ["預金：安全な保管と引き出しの自由を保証する。", "融資：預かった資金を成長期待の高い企業へ提供する。", "為替：現金を運ぶリスクを排除し、デジタル上で決済する。"])
    gen.add_message_slide("Section 02", "信用創造：銀行だけが持つ『魔法の仕組み』", 
                         ["銀行は、預かった現金以上の『預金通貨』を創り出せる。", "企業への融資が新たな預金となり、市場の通貨量が増大する。", "これが資本主義経済における成長の起爆剤となる。"])
    gen.add_diagram_slide("Section 02", "【図解】信用創造：100万円が倍以上に増えていくプロセス", gen.draw_credit_creation_diagram)
    gen.add_message_slide("Section 02", "収益モデル①：利ざや（預貸利ざや）", 
                         ["最も伝統的な収益源。資金の調達コストと運用利回りの差。", "預金者に支払う利息よりも、高い利息で企業に貸し出す。", "この『サヤ』で人件費やシステム維持費を賄う。"])
    gen.add_diagram_slide("Section 02", "【図解】利ざや：金利の『サヤ』が銀行の利益になる構造", gen.draw_interest_margin_diagram)
    gen.add_message_slide("Section 02", "収益モデル②：収益の多様化（Fee Business）", 
                         ["金利低下に対抗するため、手数料（Fee）収入を拡大。", "投資信託や保険の販売仲介手数料。M&Aや証券化のアドバイザリー。", "ソリューション提供による付加価値の獲得。"])
    gen.add_diagram_slide("Section 02", "【比較】特殊なバランスシート：預金は『負債』である", gen.draw_bs_comparison)
    gen.add_transition_slide(3, sections[2])
    gen.add_diagram_slide("Section 03", "国内銀行業界マップ：役割と規模に応じた階層構造", gen.draw_industry_map)
    gen.add_message_slide("Section 03", "マイナス金利環境の衝撃", 
                         ["長引く超低金利政策により、伝統的な利ざやが極限まで縮小。", "『預けておけば儲かる』モデルが崩壊。地方銀行を中心に危機感。", "不採算店舗の削減や、デジタルによるコスト構造改革が急務。"])
    gen.add_transition_slide(4, sections[3])
    gen.add_message_slide("Section 04", "フィンテックの台頭と決済の脱銀行化", 
                         ["スマートフォン決済の急速な普及。銀行口座が『隠れたインフラ』へ。", "プラットフォーマー（PayPay、楽天など）による顧客接点の占有。", "レンディング（融資）領域へのAI・データ活用型企業の参入。"])
    gen.add_message_slide("Section 04", "異業種参入：あらゆるサービスへの金融機能の融合", 
                         ["『銀行に行く』という行為の消失。生活のあらゆる場面に金融が溶け込む。", "スターバックスやAppleが、自社の顧客向けに便利な金融サービスを提供。", "伝統的銀行から、顧客との接点が奪われる『中抜き』の発生。"])
    gen.add_transition_slide(5, sections[4])
    gen.add_message_slide("Section 05", "DXの本質：デジタルによる価値の再定義", 
                         ["単なるネットバンキング化ではなく、UI/UXの極致を追求。", "蓄積されたビッグデータを活用した、高度な与信とコンサルティング。", "顧客のライフステージに合わせたシームレスな解決策の提案。"])
    gen.add_message_slide("Section 05", "BaaS (Banking as a Service) への転換", 
                         ["銀行自らフロントに立つのではなく、インフラ機能を外部に提供。", "他業種と連携し、金融機能を部品化（API化）して提供する。", "プラットフォームとしての生き残り戦略。"])
    gen.add_message_slide("Section 05", "2030年代の銀行像：Invisible Banking", 
                         ["『銀行』という看板が必要なくなる未来。生活に金融が溶け込む。", "AIによる完全自動の資産運用とアドバイザリー。", "社会の信頼インフラとして、データの安全性と流動性を担保する存在。"])
    gen.add_summary_slide([
        "銀行は『信用創造』という独自機能を持ち、経済を支える心臓である。",
        "伝統的な利ざやモデルから、デジタル・手数料ベースへの転換期にある。",
        "外部サービスとの金融融合により、銀行の境界線が曖昧になっている。",
        "2030年に向け、『見えない信頼インフラ』としての進化が求められる。"
    ])
    gen.add_message_slide("参考資料・サイト", "（以下、具体的なソースを追記予定）", ["全国銀行協会 統計データ", "金融庁：金融レポート", "日本銀行：金融経済統計月報"])

    output_path = "bank_strategy_presentation.pptx"
    gen.save(output_path)
    print(f"Strategic Presentation saved to {output_path}")

if __name__ == "__main__":
    generate_strategic_presentation()
