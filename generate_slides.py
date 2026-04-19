import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Constants & Design System ---
NAVY = RGBColor(0x14, 0x3D, 0x98)
LIGHT_BLUE = RGBColor(0x37, 0x5B, 0xA4)
BEIGE = RGBColor(0xE6, 0xD0, 0xA4)
OFF_WHITE = RGBColor(0xFE, 0xFE, 0xFE)
BLACK = RGBColor(0x33, 0x33, 0x33)

FONT_SANS = "Hiragino Sans"  # Preferred for macOS
FONT_SANS_BOLD = "Hiragino Sans"

class PresentationGenerator:
    def __init__(self, title, subtitle):
        self.prs = Presentation()
        # Set slide size 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle

    def _add_base_design(self, slide, title_text=None):
        """Adds standard frames and background elements inspired by Canva EAGpW_94qnc."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = OFF_WHITE

        # Top border line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = NAVY
        line.line.visible = False

        # Footer line
        footer_y = self.prs.slide_height - Inches(0.4)
        footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), footer_y, self.prs.slide_width - Inches(1.0), Inches(0.01))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = NAVY
        footer_line.line.visible = False

        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1.0))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = FONT_SANS
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.LEFT

            underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(1.5), Inches(0.05))
            underline.fill.solid()
            underline.fill.fore_color.rgb = BEIGE
            underline.line.visible = False

    def add_cover_slide(self, image_path=None):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if image_path and os.path.exists(image_path):
            # Background image
            slide.shapes.add_picture(image_path, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
            # Overlay
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width * 0.5, self.prs.slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = NAVY
            overlay.fill.transparency = 0.1
            overlay.line.visible = False
        else:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width * 0.4, self.prs.slide_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = NAVY
            rect.line.visible = False

        # Title Block - Positioned slightly higher for balance
        title_top = Inches(1.8)
        title_box = slide.shapes.add_textbox(Inches(0.5), title_top, Inches(12), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = self.title
        p.font.name = FONT_SANS
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = OFF_WHITE if image_path else NAVY
        
        p2 = tf.add_paragraph()
        p2.text = self.subtitle
        p2.font.name = FONT_SANS
        p2.font.size = Pt(20)
        p2.font.color.rgb = BEIGE
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), title_top - Inches(0.3), Inches(1.0), Inches(0.1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = BEIGE
        accent.line.visible = False

    def add_transition_slide(self, chapter_num, title):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Navy background
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.visible = False

        # Chapter Number
        num_box = slide.shapes.add_textbox(0, Inches(2.2), self.prs.slide_width, Inches(1))
        tf_num = num_box.text_frame
        p_num = tf_num.add_paragraph()
        p_num.text = f"Section {chapter_num:02d}"
        p_num.font.name = FONT_SANS
        p_num.font.size = Pt(28)
        p_num.font.color.rgb = BEIGE
        p_num.alignment = PP_ALIGN.CENTER

        # Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.prs.slide_width * 0.4, Inches(3.2), self.prs.slide_width * 0.2, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = BEIGE
        line.line.visible = False

        # Title
        title_box = slide.shapes.add_textbox(0, Inches(3.6), self.prs.slide_width, Inches(1.5))
        tf_title = title_box.text_frame
        p_title = tf_title.add_paragraph()
        p_title.text = title
        p_title.font.name = FONT_SANS
        p_title.font.size = Pt(40)
        p_title.font.bold = True
        p_title.font.color.rgb = OFF_WHITE
        p_title.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, items):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._add_base_design(slide, title)

        # Content area - Balanced margins and centered vertical anchoring
        left = Inches(0.8)
        top = Inches(1.6)
        width = self.prs.slide_width - Inches(1.6)
        height = self.prs.slide_height - Inches(2.3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE # This ensures text sits in the middle of the available height

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.name = FONT_SANS
            p.font.size = Pt(22) # Slightly increased font size for readability
            p.space_after = Pt(24) # Increased spacing between items
            p.level = 0
            if item.startswith("  "):
                p.level = 1
                p.text = item.strip()
            p.font.color.rgb = BLACK

    def save(self, filename):
        self.prs.save(filename)

def generate_banking_presentation():
    sections = [
        "1. 銀行の本質的役割と定義",
        "2. 銀行の歴史：中世から現代までの変遷",
        "3. 銀行の種類：役割に応じた多様なプレイヤー",
        "4. 3大業務と銀行収益の仕組み",
        "5. 業界を取り巻く現状と直面する課題",
        "6. 未来への展望：DXと新たな金融の形"
    ]

    gen = PresentationGenerator("銀行業界の基礎知識", "〜 経済の心臓が刻む、信頼と変革の物語 〜")
    
    # Slide 1: Cover
    gen.add_cover_slide("assets/cover_bg.png")

    # Slide 2: Table of Contents
    gen.add_content_slide("本日の目次", sections)

    # --- Section 1 ---
    gen.add_transition_slide(1, sections[0])
    gen.add_content_slide(sections[0], [
        "■ 金融（資金の融通）のプロフェッショナル",
        "  ・お金が余っている人から、足りない人へ橋渡しを行う。",
        "■ 間接金融の中核",
        "  ・預金者と借入者の間に銀行が立ち、リスクを負って資金を回す仕組み。",
        "■ 「心臓」としての役割",
        "  ・お金を血液に見立て、社会の隅々へ送り出すライフライン。"
    ])

    # --- Section 2 ---
    gen.add_transition_slide(2, sections[1])
    gen.add_content_slide(sections[1], [
        "■ ヨーロッパの起源",
        "  ・中世、金細工師（ゴールドスミス）が金を預かった「預かり証」が紙幣の原型に。",
        "  ・利息を取ることが禁忌とされていた時代、ユダヤ人が金融の先駆者となった。",
        "■ 信頼のシステム化",
        "  ・個人の信用から、組織・制度としての信用への転換。"
    ])

    gen.add_content_slide("2. 銀行の歴史（日本における成立）", [
        "■ 日本銀行の設立",
        "  ・1882年（明治15年）、中央銀行として日本銀行が誕生。",
        "  ・インフレの抑制と通貨価値の安定という重大な使命。",
        "■ 殖産興業と金融",
        "  ・明治の近代化政策を支えるため、多くの国立銀行（後の民間銀行）も設立された。"
    ])

    gen.add_content_slide("2. 銀行の歴史（高度成長〜平成再編）", [
        "■ 高度経済成長期",
        "  ・長期信用銀行（長銀）などが設備投資を支え、奇跡の成長を支えた。",
        "■ バブルの狂奔と崩壊",
        "  ・過剰融資による不良債権問題。長銀や住専の破綻と国有化。",
        "■ 平成から令和へ",
        "  ・金融ビッグバンを経て、大規模な合併（メガバンク誕生）へ。",
        "  ・異業種参入と持ち株会社制の解禁。"
    ])

    # --- Section 3 ---
    gen.add_transition_slide(3, sections[2])
    gen.add_content_slide(sections[2], [
        "■ 日本銀行（中央銀行）",
        "  ・「銀行の銀行」「政府の銀行」「発券銀行」。",
        "■ 都市銀行（メガバンク）",
        "  ・三菱UFJ、三井住友、みずほ。巨大な資本力で国・世界を支える。",
        "■ 地方銀行・第二地銀",
        "  ・「地域密着」。地元企業との信頼関係（リレーションシップ）が命。"
    ])

    gen.add_content_slide("3. 銀行の種類（デジタル・専門特化）", [
        "■ ネット銀行",
        "  ・楽天、住信SBIなど。店舗を持たず利便性と低コストを追求。",
        "■ 信託銀行",
        "  ・遺言、不動産、年金など「管理」をコアとする専門性。",
        "■ 信用金庫・信用組合",
        "  ・地域の相互扶助。利益より地域貢献を優先する非営利組織。"
    ])

    # --- Section 4 ---
    gen.add_transition_slide(4, sections[3])
    gen.add_content_slide(sections[3], [
        "1. 預金業務：大切な資産を安全に預かる。",
        "2. 貸出業務：成長を望む企業や個人へ資金を貸し付ける。",
        "   ・証書貸付、手形割引など。",
        "3. 為替業務：現金を運ばず、安全・迅速に決済を完了させる。"
    ])

    gen.add_content_slide("4. 収益の源泉：利ざやと信用創造", [
        "■ 「利ざや」の仕組み",
        "  ・貸出金利（受取）と預金利息（支払）の差額こそが利益。",
        "■ 信用創造機能",
        "  ・預金と貸出を繰り返すことで、世の中の通貨量を増大させる。",
        "  ・支払準備制度により安全性を担保しつつ、経済を膨らませる機能。"
    ])

    # --- Section 5 ---
    gen.add_transition_slide(5, sections[4])
    gen.add_content_slide(sections[4], [
        "■ 超低金利環境の長期化",
        "  ・利ざやの縮小により、旧来の預貸モデルの限界が露呈。",
        "■ プラットフォーマーとの競争",
        "  ・PayPayや楽天ポイントなど、銀行を通さない経済圏の拡大。",
        "■ 更なる代替手段の脅威",
        "  ・ステーブルコインや中央銀行デジタル通貨"
    ])

    # --- Section 6 ---
    gen.add_transition_slide(6, sections[5])
    gen.add_content_slide(sections[5], [
        "■ バンキング・アズ・ア・サービス(BaaS)",
        "  ・金融機能をパーツとして他業種に提供するプラットフォーム化。",
        "■ 融合するテクノロジー：Olive、デジタル通貨",
        "  ・銀行口座と決済アプリの完全統合。さらにCBDC（デジタル円）の可能性。",
        "■ コンサルタントへの進化",
        "  ・「金貸し」から、顧客の経営課題を解決するパートナーへ。"
    ])

    # Summary
    gen.add_content_slide("まとめ：信頼を基盤に、未来を創る", [
        "・銀行は社会の「心臓」であり、血液（資金）を循環させる責務がある。",
        "・伝統的なモデルからの脱却と、デジタルとの融合が不可避である。",
        "・「信頼」という最大の無形資産を、どうテックと掛け合わせるか。",
        "・本日の学びが、金融の未来を考えるきっかけになれば幸いです。"
    ])

    # References
    gen.add_content_slide("参考一覧", [
        "（ここに参考資料・Webサイトのリンク等を追記）"
    ])

    output_path = "bank_industry_presentation.pptx"
    gen.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    generate_banking_presentation()
